import os
import streamlit as st
import PyPDF2
from PIL import Image
import io
from google.cloud import vision
from pdf2image import convert_from_path
from pptx import Presentation
import tempfile

import vertexai
from vertexai.preview.generative_models import (
    GenerationConfig,
    GenerativeModel,
)

# Set up Google Cloud Vision Client
vision_client = vision.ImageAnnotatorClient()

# Set the page layout first
st.set_page_config(layout="wide")

# Initialize Vertex AI
PROJECT_ID = os.getenv("ADD PROJECT_ID")
LOCATION = os.getenv("ADD PROJECT LOCATION")
vertexai.init(project=PROJECT_ID, location=LOCATION)


@st.cache_resource
def load_model():
    return GenerativeModel("ADD MODEL NAME")


def create_prompt(text, option):
    if option == "Key Concepts":
        return f"""
            You are a knowledgeable subject matter expert. Identify the key concepts, terms, and their concise explanations from the text below. Format each concept and its explanation using Markdown:

            **Concept:** Explanation

            Separate each key concept with a newline.

            Input Content:
            ```{text}```

            KEY CONCEPTS:
        """
    elif option == "Short Notes":
        return f"""
            As a highly skilled subject matter expert, generate well-organized short study notes from the text below using Markdown formatting. Use bold headings for main topics and bullet points for sub-points or details.

            Example Format:
            **Main Topic 1**
            - Point 1
            - Point 2

            **Main Topic 2**
            - Sub-point A
            - Sub-point B

            Input Content:
            ```{text}```

            STUDY NOTES:
        """
    else:
        return "Invalid option selected."


def extract_text_from_image(image_file):
    try:
        with io.BytesIO(image_file.read()) as image_bytes:
            image = vision.Image(content=image_bytes.getvalue())
            response = vision_client.text_detection(image=image)
            texts = response.text_annotations
            return texts[0].description if texts else ""
    except Exception as e:
        st.error(f"Error extracting text from image: {e}")
        return ""


def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    return " ".join(page.extract_text() or "" for page in reader.pages)


def extract_text_from_pdf_images(uploaded_file):
    full_text = ""
    temp_file_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
            tmp_file.write(uploaded_file.read())
            temp_file_path = tmp_file.name

        images = convert_from_path(temp_file_path)
        for image in images:
            with io.BytesIO() as img_byte:
                image.save(img_byte, format="PNG")
                img_byte.seek(0)
                full_text += extract_text_from_image(img_byte)
    except Exception as e:
        st.error(f"Error extracting text from PDF images: {e}")
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
    return full_text


def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def get_text_response(model, prompt, config):
    responses = model.generate_content(prompt, generation_config=config, stream=True)
    return " ".join(response.text for response in responses if response.text)


# Load CSS
if os.path.exists("style.css"):
    with open("style.css") as f:
        st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)

# Streamlit app layout
st.markdown(
    """
    <head>
        <link rel="icon" href="https://freepngimg.com/save/55655-gaming-download-hd-png/512x512" sizes="16x16" type="image/png">
    </head>
    """,
    unsafe_allow_html=True,
)
st.markdown('<div class="app-container">', unsafe_allow_html=True)
st.markdown('<div class="header"><h1>StudyMate AI 📘</h1></div>', unsafe_allow_html=True)

with st.container():
    st.markdown('<div class="total-words-container"></div>', unsafe_allow_html=True) # Placeholder

    with st.sidebar:
        st.subheader("📚 About the App")
        st.write(
            "Upload a file (PDF, PPTX, TXT, Image) and generate Key Concepts or Short Notes."
        )
        uploaded_file = st.file_uploader("Upload file", type=["png", "jpg", "jpeg", "pdf", "pptx", "txt"])
        option = st.radio("Choose Output Type", options=["Key Concepts", "Short Notes"])
        generate = st.button("Generate")

    st.markdown('<div class="main-content-wrapper">', unsafe_allow_html=True)
    st.markdown('<div class="main-content">', unsafe_allow_html=True)

    if uploaded_file:
        full_text = ""
        if uploaded_file.type == "application/pdf":
            full_text = extract_text_from_pdf(uploaded_file)
            if not full_text:
                full_text = extract_text_from_pdf_images(uploaded_file)
        elif uploaded_file.type == "text/plain":
            full_text = uploaded_file.read().decode("utf-8")
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            full_text = extract_text_from_pptx(uploaded_file)
        elif uploaded_file.type in ["image/png", "image/jpeg", "image/jpg"]:
            full_text = extract_text_from_image(uploaded_file)

        word_count = len(full_text.split())
        st.markdown(f'<div class="total-words-container">Total Words Extracted: {word_count}</div>', unsafe_allow_html=True)

        if generate and full_text:
            text_snippet = full_text[:30000]
            prompt = create_prompt(text_snippet, option)
            model = load_model()
            config = {"max_output_tokens": 8000}

            with st.spinner("Generating... Please wait..."):
                response = get_text_response(model, prompt, config)
                formatted_response = response.strip()
                st.subheader(f"{option} Output:")
                st.markdown(formatted_response)

                st.session_state["generated_output"] = formatted_response
                st.session_state["original_text"] = text_snippet
                st.session_state["output_type"] = option
                st.session_state["chat_enabled"] = True
        else:
            if "generated_output" in st.session_state:
                st.subheader(f"{st.session_state['output_type']} Output:")
                st.markdown(st.session_state["generated_output"])
                st.session_state["chat_enabled"] = True

        st.markdown('</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

    if st.session_state.get("chat_enabled"):
        st.markdown("---")  # Separator
        st.subheader("💬 Ask a question about the above content")
        user_query = st.text_input("Ask a follow-up question to understand better:")
        if st.button("Get Answer") and user_query:
            chat_prompt = f"""
                You are a friendly and experienced tutor.
                A student has read the following educational content:

                ```{st.session_state['generated_output']}```

                The student has a question:
                "{user_query}"

                Please provide a very clear, beginner-friendly, and concise answer that explains the concept in a simple way.
                Do not repeat the original text. Explain it like you're teaching it to a complete beginner.
                """
            with st.spinner("Thinking..."):
                model = load_model()
                answer = get_text_response(model, chat_prompt, {"max_output_tokens": 1024})
                st.success("Answer:")
                st.write(answer)

st.markdown('</div>', unsafe_allow_html=True) 
